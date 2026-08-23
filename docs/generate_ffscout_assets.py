# -*- coding: utf-8 -*-
"""Generate FFScout CEO deck (PPTX) and FPL manager survey (XLSX)."""

from pathlib import Path

from openpyxl import Workbook
from openpyxl.styles import Alignment, Font, PatternFill, Border, Side
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

OUT_DIR = Path(__file__).resolve().parent

# Palette — flat, professional, no purple-on-white AI cliché
NAVY = RGBColor(0x0F, 0x1C, 0x2E)
INK = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x5C, 0x67, 0x73)
ACCENT = RGBColor(0x0B, 0x6E, 0x4F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF4, 0xF6, 0xF8)
LINE = RGBColor(0xD0, 0xD7, 0xDE)


def set_run(run, text, size=18, bold=False, color=INK, font_name="Calibri"):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name


def add_textbox(slide, left, top, width, height, text, size=18, bold=False, color=INK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    set_run(p.add_run() if p.runs else p.runs[0] if False else None, text, size, bold, color)
    # pptx paragraphs always have one empty run path — do cleanly:
    p.clear() if hasattr(p, "clear") else None
    # Rebuild paragraph properly
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    set_run(run, text, size, bold, color)
    return box


def clear_and_write(tf, lines, size=16, bold=False, color=INK, space_after=8):
    """lines: list of str, or list of (str, bold, size) tuples."""
    tf.clear()
    tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, tuple):
            text, is_bold, sz = item[0], item[1] if len(item) > 1 else bold, item[2] if len(item) > 2 else size
            col = item[3] if len(item) > 3 else color
        else:
            text, is_bold, sz, col = item, bold, size, color
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(space_after)
        run = p.add_run()
        set_run(run, text, sz, is_bold, col)


def add_bullets(slide, left, top, width, height, bullets, size=15):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(6)
        run = p.add_run()
        set_run(run, f"•  {b}", size, False, INK)
    return box


def add_footer(slide, page, total):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(8.5), Inches(0.3))
    tf = box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    set_run(run, f"FALEAGUE AI  ·  Confidential conversation material  ·  {page}/{total}", 10, False, MUTED)
    box2 = slide.shapes.add_textbox(Inches(8.8), Inches(7.05), Inches(0.8), Inches(0.3))
    p2 = box2.text_frame.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    r2 = p2.add_run()
    set_run(r2, "", 10, False, MUTED)


def add_accent_bar(slide):
    shape = slide.shapes.add_shape(
        1,  # rectangle
        Inches(0),
        Inches(0),
        Inches(0.12),
        Inches(7.5),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()


def blank_slide(prs):
    blank = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(blank)
    # light bg via full rect
    bg = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()
    # send to back is automatic for first shape roughly — keep accent
    add_accent_bar(slide)
    return slide


def title_block(slide, title, subtitle=None):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(12.2), Inches(0.7))
    tf = box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    set_run(run, title, 28, True, NAVY)
    if subtitle:
        box2 = slide.shapes.add_textbox(Inches(0.55), Inches(0.95), Inches(12.2), Inches(0.4))
        p2 = box2.text_frame.paragraphs[0]
        r2 = p2.add_run()
        set_run(r2, subtitle, 14, False, MUTED)


def build_pptx():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    total = 14

    # --- Title ---
    slide = blank_slide(prs)
    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.5), Inches(1.2))
    tf = box.text_frame
    p = tf.paragraphs[0]
    set_run(p.add_run(), "FPL in China × Fantasy Football Scout", 36, True, NAVY)
    box2 = slide.shapes.add_textbox(Inches(0.8), Inches(3.3), Inches(11.5), Inches(0.5))
    set_run(box2.text_frame.paragraphs[0].add_run(), "Collaboration conversation — FALEAGUE AI", 20, False, ACCENT)
    add_bullets(
        slide,
        Inches(0.8),
        Inches(4.2),
        Inches(10),
        Inches(2),
        [
            "Presenter: Ray Wong · FALEAGUE AI / Faleague",
            "Product: faleague-ai.com  ·  Sister: faleague.cn",
            "Tone: complementary · evidence-led · no vanity metrics",
        ],
        size=16,
    )
    add_footer(slide, 1, total)

    # --- Agenda / how to use ---
    slide = blank_slide(prs)
    title_block(slide, "How we’ll use this conversation", "≈20–25 min + discussion")
    add_bullets(
        slide,
        Inches(0.7),
        Inches(1.6),
        Inches(11.5),
        Inches(4.5),
        [
            "Lead with mutual fit — not a product tour",
            "Be explicit: what is built & live vs hypothesised",
            "No invented China FPL headcount",
            "This is v1 — after the survey, we revise into a v2 evidence deck",
            "Ask: comfort with survey-first, then return with proof — not a big-bang JV",
            "Leave half the call for your questions",
        ],
        size=18,
    )
    add_footer(slide, 2, total)

    # --- Slide 1 why ---
    slide = blank_slide(prs)
    title_block(slide, "1 · Why we’re talking", "Opening the China FPL opportunity — together")
    add_bullets(
        slide,
        Inches(0.7),
        Inches(1.6),
        Inches(11.5),
        Inches(3.2),
        [
            "Fantasy Football Scout is the category authority for serious FPL managers",
            "We built a live bilingual FPL product with China-native distribution (中文 + WeChat)",
            "This is not “please promote us” — it’s: does a China channel create value for Scout?",
            "What would a light, reversible collaboration look like?",
        ],
        size=17,
    )
    box = slide.shapes.add_shape(1, Inches(0.7), Inches(5.2), Inches(11.8), Inches(1.1))
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT
    box.line.color.rgb = LINE
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    set_run(
        p.add_run(),
        "One line: Scout owns trust and depth. We own China access and product distribution. The overlap is the opportunity.",
        16,
        True,
        NAVY,
    )
    add_footer(slide, 3, total)

    # --- Market ---
    slide = blank_slide(prs)
    title_block(slide, "2 · The market (realistic framing)")
    add_bullets(
        slide,
        Inches(0.7),
        Inches(1.5),
        Inches(5.8),
        Inches(4.5),
        [
            "Global FPL ~11.5M managers in 2024/25 (order of magnitude)",
            "English content saturates the serious segment",
            "China fantasy sports industry is large & growing — this is NOT FPL China revenue",
            "FPL-specific China headcount is under-published — we will not invent a number",
        ],
        size=15,
    )
    add_bullets(
        slide,
        Inches(6.8),
        Inches(1.5),
        Inches(5.8),
        Inches(4.5),
        [
            "What we know operationally: Chinese managers need language, timezone, and WeChat-native habits",
            "UK Twitter / Discord tools do not fully serve that habit loop",
            "Honest gap: trusted Chinese-language research + tools + distribution is thin vs the English ecosystem Scout leads",
        ],
        size=15,
    )
    add_footer(slide, 4, total)

    # --- China gap ---
    slide = blank_slide(prs)
    title_block(slide, "3 · The China gap", "Observed behaviour — say this clearly")
    add_bullets(
        slide,
        Inches(0.7),
        Inches(1.45),
        Inches(12),
        Inches(4.8),
        [
            "Chinese FPLers commonly pull recommendations from overseas FPL creators & sites",
            "Local creators / fans often translate or rewrite that content onto 小红书 (XHS), WeChat, Douyin, Bilibili",
            "China already consumes Scout-ecosystem signal — often with informal translation and weak attribution",
            "Strategic question: stay unofficial & uncontrolled, or become an attributed, quality-controlled China channel?",
            "Entering China cold is expensive; partnering is lower risk — and a brand-protection move",
            "Need: Chinese GW prep · bilingual tools · WeChat/XHS habit loop · a partner who won’t dilute quality",
        ],
        size=15,
    )
    add_footer(slide, 5, total)

    # --- Who we are ---
    slide = blank_slide(prs)
    title_block(slide, "4 · Who we are / what we built", "FALEAGUE AI — AI analyst for FPL, expanding into a football hub")
    add_bullets(
        slide,
        Inches(0.7),
        Inches(1.5),
        Inches(12),
        Inches(4.8),
        [
            "Live: https://www.faleague-ai.com (EN + ZH)  ·  Sister: faleague.cn",
            "Real stack: Next.js + Supabase + Gemini tool-use · deterministic TS optimisers + LLM narration",
            "FPL Insights (free + Insights Pro path) · World Cup hub · news / daily briefing · Mini 5",
            "WeChat 小程序 shell + automated daily Chinese cards",
            "Credibility: our WeChat digests already cite Fantasy Football Scout — amplifying signal into Chinese feeds",
        ],
        size=16,
    )
    add_footer(slide, 6, total)

    # --- Demo ---
    slide = blank_slide(prs)
    title_block(slide, "5 · Product snapshot", "Show live — 3 minutes, not more slides")
    add_bullets(
        slide,
        Inches(0.7),
        Inches(1.5),
        Inches(12),
        Inches(4.8),
        [
            "/zh/fpl/insights — research tables from live FPL data",
            "Free live: Pre-season · Best of Position · Set-pieces · DefCon · Fixture swing",
            "Premium path: Transfers · Differentials · Prices · xG/xA gap · xP accuracy",
            "/zh/chat — AI grounded in tools + Entry ID (not freeform hallucination)",
            "WeChat artefact — daily card / Best-of-Position hook with deep links",
            "Positioning: we are NOT claiming Opta parity — China-facing product + distribution layer",
        ],
        size=15,
    )
    add_footer(slide, 7, total)

    # --- Marketing ---
    slide = blank_slide(prs)
    title_block(slide, "6 · Marketing plan & focus", "North star: Chinese-speaking FPL habit loop → light monetisation")
    cols = [
        (
            "Focus 1 · WeChat-first",
            [
                "Daily Chinese digests",
                "Best of Position hooks → Insights",
                "GW1 drafts in Beijing time",
                "Mini program re-entry",
            ],
        ),
        (
            "Focus 2 · Product retention",
            [
                "Entry ID → planner / captain / transfers",
                "Insights = research desk",
                "AI chat = analyst",
                "Free core strong; Pro for depth",
            ],
        ),
        (
            "Focus 3 · Content flywheel",
            [
                "Authorised attributed packaging",
                "Meet managers on WeChat + XHS",
                "Protect Scout brand equity",
                "Measure CTR / binds / returns",
            ],
        ),
    ]
    for i, (h, items) in enumerate(cols):
        left = Inches(0.55 + i * 4.2)
        hb = slide.shapes.add_textbox(left, Inches(1.5), Inches(4.0), Inches(0.45))
        set_run(hb.text_frame.paragraphs[0].add_run(), h, 16, True, ACCENT)
        add_bullets(slide, left, Inches(2.05), Inches(4.0), Inches(3.2), items, size=14)
    note = slide.shapes.add_textbox(Inches(0.7), Inches(5.5), Inches(11.5), Inches(1.0))
    clear_and_write(
        note.text_frame,
        [
            ("Not doing: compete for Scout UK membership · invent China user millions · heavy ads before WeChat retention proves out", False, 13, MUTED),
        ],
    )
    add_footer(slide, 8, total)

    # --- Business ---
    slide = blank_slide(prs)
    title_block(slide, "7 · Business potential (China × tools)", "Keep this sober")
    rows = [
        ("Layer", "Near-term (this season)", "Medium-term"),
        ("Audience", "Chinese FPL managers via WeChat + faleague.cn", "Broader Greater China PL fans entering FPL"),
        ("Product", "Free Insights + AI hub", "Insights Pro + sponsored Insights slots"),
        ("Brand", "Faleague as China FPL hub", "Co-branded Scout × Faleague trust signal"),
        ("Scout upside", "Incremental reach + brand presence in CN", "Membership / licensing / co-products"),
    ]
    top = Inches(1.55)
    for r, row in enumerate(rows):
        y = top + Inches(r * 0.55)
        for c, cell in enumerate(row):
            x = Inches(0.55 + c * 4.15)
            w = Inches(4.05)
            shape = slide.shapes.add_shape(1, x, y, w, Inches(0.5))
            shape.fill.solid()
            shape.fill.fore_color.rgb = NAVY if r == 0 else (LIGHT if r % 2 else WHITE)
            shape.line.color.rgb = LINE
            tf = shape.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            set_run(p.add_run(), cell, 12, r == 0, WHITE if r == 0 else INK)
    note = slide.shapes.add_textbox(Inches(0.7), Inches(4.5), Inches(11.5), Inches(1.5))
    clear_and_write(
        note.text_frame,
        [
            ("Monetisation honesty", True, 14, ACCENT),
            ("Billing plumbing exists; premium enforcement is soft while habit grows. China may also mean sponsorship + membership referral — not only SaaS.", False, 14, INK),
        ],
    )
    add_footer(slide, 9, total)

    # --- Collaboration ---
    slide = blank_slide(prs)
    title_block(slide, "8 · Collaboration options", "Present as a menu — pick 1–2 to start")
    add_bullets(
        slide,
        Inches(0.7),
        Inches(1.45),
        Inches(12),
        Inches(4.5),
        [
            "A. Content & attribution (lightest) — Scout-approved Chinese WeChat/XHS cards with clear source + link",
            "B. Co-branded China Insights — Scout editorial + our bilingual UI; Scout keeps brand control",
            "C. Membership funnel experiment — soft CTA to Scout membership; unique links",
            "D. Joint research — bilingual unbranded survey → revise this deck to v2 with evidence → optional co-branded report",
            "Recommended: D first (or A + D) — survey → evidence deck → then scale content pilot",
        ],
        size=15,
    )
    add_footer(slide, 10, total)

    # --- Learnings ---
    slide = blank_slide(prs)
    title_block(slide, "9 · What we’ve learned building this", "Founder credibility — not tech vanity")
    add_bullets(
        slide,
        Inches(0.7),
        Inches(1.5),
        Inches(12),
        Inches(4.8),
        [
            "Optimisers first, LLM second — rankings must be deterministic; AI explains",
            "China is a distribution problem before a model problem",
            "Trust travels with attribution — packaging + language is the unlock",
            "Hub > single feature — Insights + chat + Entry ID + daily card",
            "We’re early on measured growth — product live; survey + pilot produce proof",
        ],
        size=17,
    )
    add_footer(slide, 11, total)

    # --- 90 day ---
    slide = blank_slide(prs)
    title_block(slide, "10 · Proposed 90-day pilot")
    pilot = [
        ("Week 1–3", "Field bilingual unbranded survey (pilot 30–50 → 200–500+)", "Clean dataset"),
        ("Week 3–5", "Analyse China vs RoW; revise this deck to v2 with survey evidence", "Evidence pack"),
        ("Week 5–6", "Share v2 with Scout; align brand/attribution if interest holds", "Written one-pager"),
        ("Week 6–9", "Optional: WeChat/XHS pilot cards with Scout attribution + unique links", "CTR / deep-link opens"),
        ("Week 9–12", "Go / no-go on co-branded Insights asset + membership interest", "Qualitative + link metrics"),
    ]
    for i, (w, a, s) in enumerate(pilot):
        y = Inches(1.45 + i * 0.7)
        shape = slide.shapes.add_shape(1, Inches(0.55), y, Inches(12.2), Inches(0.62))
        shape.fill.solid()
        shape.fill.fore_color.rgb = LIGHT if i % 2 == 0 else WHITE
        shape.line.color.rgb = LINE
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        set_run(p.add_run(), f"{w}  ·  {a}  →  {s}", 13, False, INK)
    ask = slide.shapes.add_textbox(Inches(0.7), Inches(5.2), Inches(11.5), Inches(1.0))
    clear_and_write(
        ask.text_frame,
        [
            ("Ask (v1 call): comfort with us fielding the unbranded survey and returning with a v2 evidence deck — plus a point of contact for early sight of the instrument.", True, 14, NAVY),
        ],
    )
    add_footer(slide, 12, total)

    # --- Discussion ---
    slide = blank_slide(prs)
    title_block(slide, "11 · Discussion prompts")
    add_bullets(
        slide,
        Inches(0.7),
        Inches(1.6),
        Inches(12),
        Inches(4.8),
        [
            "Is China a strategic priority for Scout in the next 12–24 months, or opportunistic?",
            "What must never be localised / scraped / paraphrased without editorial control?",
            "Prefer brand reach only, or also a commercial path (membership / licensing)?",
            "Would a joint survey / report be useful to Scout’s audience as well?",
            "Who would own the pilot internally if we proceed?",
        ],
        size=18,
    )
    add_footer(slide, 13, total)

    # --- Close ---
    slide = blank_slide(prs)
    title_block(slide, "12 · Closing")
    add_bullets(
        slide,
        Inches(0.7),
        Inches(1.55),
        Inches(12),
        Inches(3.2),
        [
            "Scout = trust + depth in FPL",
            "FALEAGUE = live bilingual product + WeChat China distribution",
            "Ask = reversible 90-day content + research pilot — not a big-bang JV",
        ],
        size=18,
    )
    leave = slide.shapes.add_textbox(Inches(0.7), Inches(4.6), Inches(11.5), Inches(1.6))
    clear_and_write(
        leave.text_frame,
        [
            ("Leave-behind", True, 16, ACCENT),
            ("1) This deck (v1)   2) Live demo links   3) Unbranded FPLer survey (Excel — ¥ pricing)   → then v2 evidence deck after results", False, 14, INK),
            ("Intro line: I’ve built a live bilingual FPL product with WeChat distribution — I’d like to field an unbranded survey and return with evidence for a small attributed pilot.", False, 13, MUTED),
        ],
        space_after=6,
    )
    add_footer(slide, 14, total)

    out = OUT_DIR / "FFScout-China-Collaboration-Deck-v1.pptx"
    prs.save(out)
    return out


def style_header(ws, row, cols):
    fill = PatternFill("solid", fgColor="0F1C2E")
    font = Font(bold=True, color="FFFFFF", name="Calibri", size=11)
    for col in range(1, cols + 1):
        cell = ws.cell(row=row, column=col)
        cell.fill = fill
        cell.font = font
        cell.alignment = Alignment(wrap_text=True, vertical="center")


def autosize(ws, max_width=48):
    for col in ws.columns:
        letter = get_column_letter(col[0].column)
        length = 0
        for cell in col:
            if cell.value:
                length = max(length, min(len(str(cell.value)), max_width))
        ws.column_dimensions[letter].width = max(12, length + 2)


def build_xlsx():
    wb = Workbook()

    # --- Instructions ---
    ws = wb.active
    ws.title = "00_Instructions"
    instructions = [
        ["FPL Manager Survey — fielding workbook"],
        ["Bias rule: Do NOT name Fantasy Football Scout (or other partner brands) in any respondent-facing copy."],
        ["Languages: English + 简体中文 share the same Question ID."],
        ["Use sheets EN_Questions / ZH_Questions to import into Wenjuanxing / Typeform / Google Forms."],
        ["Options sheet lists answer choices; one row per option."],
        ["Skip_Logic sheet documents conditional display."],
        ["Analysis_Plan is internal only — not for respondents."],
        ["Target: pilot 30–50, then scale toward 200–500+; report China vs Rest of World separately."],
        ["Incentive: Faleague Insights Pro / league prize / generic premium trial — no named third-party membership."],
    ]
    for r, row in enumerate(instructions, 1):
        ws.cell(r, 1, row[0])
        ws.cell(r, 1).font = Font(bold=(r == 1), size=14 if r == 1 else 11, name="Calibri")
        ws.cell(r, 1).alignment = Alignment(wrap_text=True)
    ws.column_dimensions["A"].width = 110

    # Question bank — structured
    # id, section, type, en, zh, max_select, skip_note
    questions = [
        ("S1", "Screener", "single", "Do you currently play Fantasy Premier League (official FPL)?", "你目前是否在玩官方 Fantasy Premier League（FPL）？", "", ""),
        ("S2", "Screener", "single", "Where are you primarily based?", "你主要居住地？", "", ""),
        ("S3", "Screener", "single", "Preferred language for FPL content", "你更希望用哪种语言看 FPL 内容？", "", ""),
        ("Q1", "A_Profile", "single", "How many full FPL seasons have you played?", "完整打过几个赛季？", "", ""),
        ("Q2", "A_Profile", "single", "Typical overall rank band (best finish or current feel)", "大致排名带？", "", ""),
        ("Q3", "A_Profile", "multi", "Main motivation (pick up to 2)", "主要动机（最多2项）", "2", ""),
        ("Q4", "A_Profile", "single", "Hours spent on FPL in a typical Gameweek", "每周投入时间", "", ""),
        ("Q5", "B_Sources", "multi", "Which source types do you use regularly for FPL decisions?", "常用来做决定的信息源类型（多选）", "", "No brand names in options"),
        ("Q6", "B_Sources", "rank3", "Rank your top 3 source types from Q5 (ordered)", "从Q5中选出最重要的前3个并排序", "3", ""),
        ("Q7", "B_Sources", "open", "Name up to 3 FPL sites, tools, or creators you trust most for decisions (unaided)", "请写出你最信任的最多3个FPL网站/工具/创作者名称（未提示）", "", "Do not prompt brands"),
        ("Q8", "B_Sources", "multi", "How do you usually get value from overseas FPL content?", "你通常如何使用海外FPL内容？", "", ""),
        ("Q9", "B_Sources", "single", "How much do you trust English-language FPL analysis if you prefer Chinese?", "（偏中文用户）对英文FPL分析的依赖程度", "", "Skip if S3 = English only"),
        ("Q10", "C_Tools", "multi", "Which decisions are hardest each GW? (pick up to 3)", "每周最难的决定（最多3）", "3", ""),
        ("Q11", "C_Tools", "multi", "Which tool types would you pay for?", "愿意付费的工具类型（多选）", "", ""),
        ("Q12", "C_Tools", "single", "Willing to pay for premium FPL tools / membership? (prices in CNY; ≈GBP for reference)", "可接受月费（人民币）", "", ""),
        ("Q13", "C_Tools", "single", "Preferred device for FPL research", "主要设备", "", ""),
        ("Q14", "D_China", "multi", "Where do you usually discuss or browse FPL with others?", "平时在哪聊或刷 FPL？（多选）", "", "Show if China/HK/TW/SG-MY or Chinese/bilingual"),
        ("Q15", "D_China", "single", "In the last month, how often did you see Chinese posts that clearly came from overseas FPL tips / analysis (translated or rewritten)?", "近一个月，你看到「明显来自海外FPL建议/分析、被翻译或改写」的中文帖频率？", "", "China module"),
        ("Q16", "D_China", "single", "When you see those translated / rewritten posts, what do you usually do?", "看到这类翻译/改写帖时，你通常？", "", "China module"),
        ("Q17", "D_China", "multi", "Ideal Chinese FPL content format (pick up to 2)", "最想要的中文内容形态（最多2）", "2", "China module"),
        ("Q18", "D_China", "single", "Biggest barrier to using overseas FPL sites/tools directly", "直接使用海外FPL站点/工具的最大障碍", "", "China module"),
        ("Q19", "D_China", "single", "Interest in a trusted Chinese hub that packages high-quality overseas-grade analysis + tools (with clear source attribution)", "对「可信的中文FPL中心（海外级分析质量+工具，且标明来源）」兴趣", "", "China module"),
        ("Q20", "E_Faleague", "single", "Have you used faleague-ai.com / FALEAGUE AI?", "是否用过 faleague-ai.com？", "", "Optional module"),
        ("Q21", "E_Faleague", "multi", "Which parts matter most?", "（若用过）最有用的部分（多选）", "", "If Q20 used"),
        ("Q22", "E_Faleague", "scale0_10", "How likely are you to recommend FALEAGUE to a friend who plays FPL? (0–10)", "推荐意愿 0–10", "", "If Q20 used"),
        ("Q23", "E_Faleague", "open", "What’s missing that would make you return every GW?", "缺什么会让你每周回来？", "", "If Q20 used"),
        ("Q24", "F_Trust", "single", "If a well-known overseas FPL research brand offered Chinese-localised content or tools via a China partner, would you…", "若某家知名海外FPL研究品牌通过中国合作伙伴提供中文内容/工具，你会？", "", "Do not name brands"),
        ("Q25", "F_Trust", "multi", "What would make that feel trustworthy?", "怎样才显得可信？（多选）", "", ""),
        ("Q26", "F_Trust", "open", "One sentence — what would most improve your FPL season?", "一句话：最能提升你本赛季体验的是？", "", ""),
        ("D1", "Demographics", "single", "Age band", "年龄段", "", "Optional"),
        ("D2", "Demographics", "single", "Gender", "性别", "", "Optional"),
        ("D3", "Demographics", "open", "Email / WeChat for prize draw only (optional)", "抽奖用邮箱/微信（可选）", "", "Optional"),
    ]

    options = {
        "S1": [
            ("Yes, this season", "是，本赛季"),
            ("Yes, but not this season yet", "是，但本赛季还没开始"),
            ("I used to", "以前玩过"),
            ("No", "否"),
        ],
        "S2": [
            ("Mainland China", "中国大陆"),
            ("Hong Kong / Macau / Taiwan", "港澳台"),
            ("Singapore / Malaysia", "新马"),
            ("United Kingdom", "英国"),
            ("Rest of Europe", "欧洲其他"),
            ("North America", "北美"),
            ("India / South Asia", "印度南亚"),
            ("Africa / Middle East", "中东非洲"),
            ("Other (specify)", "其他"),
        ],
        "S3": [
            ("English only", "仅英文"),
            ("Chinese only", "仅中文"),
            ("Both (bilingual)", "中英双语"),
            ("Other", "其他"),
        ],
        "Q1": [("First season", "首季"), ("2–3", "2–3"), ("4–6", "4–6"), ("7+", "7+")],
        "Q2": [
            ("Top 10k", "前1万"),
            ("Top 100k", "前10万"),
            ("Top 1M", "前100万"),
            ("Outside 1M", "100万以外"),
            ("Prefer not to say", "不愿说"),
        ],
        "Q3": [
            ("Beat friends / mini-league", "打败朋友小联赛"),
            ("Climb overall rank", "冲总榜"),
            ("Content / creator community", "看内容社区"),
            ("Learn football tactics / players", "学球"),
            ("Casual fun only", "随便玩玩"),
        ],
        "Q4": [("<1h", "<1小时"), ("1–3h", "1–3"), ("3–5h", "3–5"), ("5h+", "5+")],
        "Q5": [
            ("Official FPL app / site", "官方FPL"),
            ("Overseas FPL websites / membership tools", "海外FPL网站或会员工具"),
            ("Overseas FPL content creators (YouTube, podcasts, X/Twitter, etc.)", "海外FPL内容创作者（YouTube、播客、X等）"),
            ("Chinese FPL creators / 公众号 / KOLs", "中文FPL创作者·公众号·KOL"),
            ("Local posts that translate or summarise overseas FPL content (XHS / WeChat / Douyin / Bilibili)", "把海外内容翻译或改写成中文的本地帖（小红书/微信/抖音/B站等）"),
            ("WeChat groups / Moments", "微信群·朋友圈"),
            ("小红书 (Xiaohongshu)", "小红书"),
            ("Douyin / Bilibili / other short video", "抖音·B站等短视频"),
            ("Reddit / Discord / Telegram", "Reddit·Discord·Telegram"),
            ("Friends / mini-league chat only", "只问朋友"),
            ("AI chat tools (ChatGPT, etc.)", "AI聊天工具"),
            ("Other (specify)", "其他"),
        ],
        "Q8": [
            ("I consume it directly in English", "直接看英文原文"),
            ("I use translation tools myself", "自己用翻译工具"),
            ("I rely on Chinese creators / posts that translate or rewrite it", "依赖中文创作者或翻译转载帖"),
            ("Friends summarise it for me", "朋友帮忙总结"),
            ("I rarely use overseas content", "很少用海外内容"),
        ],
        "Q9": [
            ("I rely on it heavily", "重度依赖"),
            ("I use it when translated / summarised", "需要翻译摘要才用"),
            ("I rarely use it (language barrier)", "语言障碍很少用"),
            ("I don’t use it", "不用"),
        ],
        "Q10": [
            ("Captain", "队长"),
            ("Transfers (hits vs hold)", "转会"),
            ("Differentials", "冷门"),
            ("Bench / order", "板凳"),
            ("Chip timing", "Chip"),
            ("Price changes", "涨跌价"),
            ("Fixture planning", "赛程"),
            ("Set-piece / minutes certainty", "定位球与出场时间"),
        ],
        "Q11": [
            ("Projected points / xP", "xP投影"),
            ("Predicted lineups", "预测首发"),
            ("Transfer planner", "转会规划"),
            ("Fixture ticker / FDR", "赛程FDR"),
            ("Price change predictor", "涨跌价预测"),
            ("Differentials / ownership", "冷门与拥有率"),
            ("Team rating / RMT-style review", "阵容评分"),
            ("AI chat analyst tied to my squad", "绑定我阵容的AI分析"),
            ("Chinese-language daily digest", "中文每日简报"),
            ("None — free only", "只用免费"),
        ],
        "Q12": [
            ("No", "不愿付"),
            ("Yes, under ¥20 / month (≈ £2)", "¥20以下/月"),
            ("¥20–40 / month (≈ £2–4)", "¥20–40/月"),
            ("¥40–70 / month (≈ £4–7)", "¥40–70/月"),
            ("¥70–100 / month (≈ £7–10)", "¥70–100/月"),
            ("¥100+ / month (≈ £10+)", "¥100+/月"),
            ("Prefer annual only", "更想买年费"),
        ],
        "Q13": [
            ("Mobile phone primarily", "手机"),
            ("Desktop primarily", "电脑"),
            ("Both equally", "都用"),
        ],
        "Q14": [
            ("WeChat group", "微信群"),
            ("WeChat Moments", "朋友圈"),
            ("公众号", "公众号"),
            ("Mini program", "小程序"),
            ("QQ", "QQ"),
            ("小红书 (Xiaohongshu)", "小红书"),
            ("Douyin / Bilibili", "抖音B站"),
            ("English platforms (X, Discord, Reddit)", "英文平台"),
            ("Offline friends only", "线下朋友"),
            ("I don’t discuss", "不聊"),
        ],
        "Q15": [
            ("Almost every day", "几乎每天"),
            ("A few times a week", "一周几次"),
            ("Occasionally", "偶尔"),
            ("Rarely / never", "很少或没有"),
            ("Not sure", "不确定"),
        ],
        "Q16": [
            ("Treat them as a main input for my team", "当作主要依据"),
            ("Use them as one input among several", "当作参考之一"),
            ("Check the original overseas source when I can", "有机会会去核对原文"),
            ("Mostly ignore them", "基本忽略"),
            ("I don’t see this kind of content", "没见过这类内容"),
        ],
        "Q17": [
            ("Short daily text card", "短日报卡片"),
            ("Long analysis article", "长文分析"),
            ("Short video", "短视频"),
            ("Live stream / Q&A", "直播问答"),
            ("Interactive tools / tables", "互动数据表"),
            ("AI Q&A", "AI问答"),
        ],
        "Q18": [
            ("Language", "语言"),
            ("Paywall / price", "付费墙"),
            ("VPN / access friction", "网络访问"),
            ("Don’t know good ones", "不知道好的"),
            ("Prefer local creators / translated feeds", "更信本地创作者或翻译帖"),
            ("No barrier", "无障碍"),
        ],
        "Q19": [
            ("Very interested", "很感兴趣"),
            ("Somewhat", "有点兴趣"),
            ("Neutral", "一般"),
            ("Not interested", "不感兴趣"),
        ],
        "Q20": [
            ("Yes, regularly", "经常"),
            ("Yes, once or twice", "一两回"),
            ("Heard of", "听说过"),
            ("No", "没有"),
        ],
        "Q21": [
            ("AI chat", "AI聊天"),
            ("Insights / Best of Position", "Insights位置精选"),
            ("Pre-season signals", "季前信号"),
            ("Planner / transfers", "规划转会"),
            ("Fixtures / FDR", "赛程"),
            ("WeChat daily card", "微信日报"),
            ("World Cup hub", "世界杯"),
            ("Mini 5", "Mini5"),
            ("News / FPL daily", "新闻"),
        ],
        "Q24": [
            ("Follow / subscribe", "关注订阅"),
            ("Consider paid membership", "考虑付费会员"),
            ("Use free tools only", "只用免费"),
            ("Not interested", "不感兴趣"),
            ("Need to see quality first", "先看质量"),
        ],
        "Q25": [
            ("Clear original-source attribution", "明确原文署名"),
            ("Accurate translations (not rough machine-only)", "高质量翻译"),
            ("Same analytical depth as the overseas original", "分析深度不输海外原文"),
            ("Local community moderators", "本地社群管理"),
            ("Fair China pricing", "合理中国定价"),
            ("No gambling framing", "无博彩包装"),
            ("Consistent GW cadence", "稳定的每周更新"),
        ],
        "D1": [("<18", "<18"), ("18–24", "18–24"), ("25–34", "25–34"), ("35–44", "35–44"), ("45+", "45+"), ("Prefer not", "不愿说")],
        "D2": [("Male", "男"), ("Female", "女"), ("Other / prefer not", "其他/不愿说")],
    }

    # EN Questions sheet
    ws_en = wb.create_sheet("EN_Questions")
    headers = ["Question_ID", "Section", "Type", "Question_EN", "Max_Select", "Notes"]
    ws_en.append(headers)
    style_header(ws_en, 1, len(headers))
    for q in questions:
        ws_en.append([q[0], q[1], q[2], q[3], q[5], q[6]])
    for row in ws_en.iter_rows(min_row=2, max_row=ws_en.max_row, min_col=1, max_col=6):
        for cell in row:
            cell.alignment = Alignment(wrap_text=True, vertical="top")
            cell.font = Font(name="Calibri", size=11)
    autosize(ws_en)

    # ZH Questions
    ws_zh = wb.create_sheet("ZH_Questions")
    headers_zh = ["Question_ID", "Section", "Type", "Question_ZH", "Max_Select", "Notes"]
    ws_zh.append(headers_zh)
    style_header(ws_zh, 1, len(headers_zh))
    for q in questions:
        ws_zh.append([q[0], q[1], q[2], q[4], q[5], q[6]])
    for row in ws_zh.iter_rows(min_row=2, max_row=ws_zh.max_row, min_col=1, max_col=6):
        for cell in row:
            cell.alignment = Alignment(wrap_text=True, vertical="top")
            cell.font = Font(name="Calibri", size=11)
    autosize(ws_zh)

    # Options
    ws_opt = wb.create_sheet("Options")
    ws_opt.append(["Question_ID", "Option_Order", "Option_EN", "Option_ZH"])
    style_header(ws_opt, 1, 4)
    for qid, opts in options.items():
        for i, (en, zh) in enumerate(opts, 1):
            ws_opt.append([qid, i, en, zh])
    autosize(ws_opt)

    # Skip logic
    ws_skip = wb.create_sheet("Skip_Logic")
    ws_skip.append(["Question_ID", "Show_When", "Hide_When", "Notes"])
    style_header(ws_skip, 1, 4)
    skip_rows = [
        ("S1", "Always", "If No → end or soft PL fan exit", "Screener"),
        ("Q9", "S3 = Chinese only OR Both", "S3 = English only", "Language trust"),
        ("Q14–Q19", "S2 in Mainland China / HK-Macau-TW / SG-MY OR S3 Chinese/bilingual", "Else", "China social module"),
        ("Q21–Q23", "Q20 = Yes regularly OR once/twice", "Q20 = Heard of / No", "FALEAGUE module"),
        ("Q7", "Always (optional open)", "", "Unaided — no brand prompts"),
        ("Q24–Q25", "Always", "", "Keep unbranded wording"),
    ]
    for row in skip_rows:
        ws_skip.append(list(row))
    autosize(ws_skip)

    # Analysis plan
    ws_an = wb.create_sheet("Analysis_Plan")
    ws_an.append(["Order", "Analysis block", "Key questions", "Claim discipline"])
    style_header(ws_an, 1, 4)
    analysis = [
        (1, "Sample composition", "S2, S3, Q1, Q2", "Report China vs RoW separately"),
        (2, "Source map", "Q5, Q6", "Category-level; no partner brand in fielding"),
        (3, "Translation bridge", "Q8, Q15, Q16", "Core China insight for deck narrative"),
        (4, "Unaided trust list", "Q7", "Code open answers AFTER fieldwork"),
        (5, "Jobs to be done × pay", "Q10, Q11, Q12", "Product & membership packaging"),
        (6, "Packaging / hub signal", "Q19, Q24, Q25", "Unbranded partnership signal"),
        (7, "FALEAGUE product", "Q20–Q23", "Only if n allows"),
    ]
    for row in analysis:
        ws_an.append(list(row))
    ws_an.append([])
    ws_an.append(["Safe claim shape:", "Among surveyed Chinese-speaking managers, Y% regularly use overseas creators/sites, and Z% rely on local translated/rewritten posts on XHS/WeChat."])
    ws_an.append(["Avoid until unaided coding:", "Any single overseas brand’s awareness % from this wave."])
    autosize(ws_an)

    # Response template (empty collection sheet)
    ws_resp = wb.create_sheet("Response_Template")
    resp_headers = ["Response_ID", "Timestamp", "Channel", "Language"] + [q[0] for q in questions]
    ws_resp.append(resp_headers)
    style_header(ws_resp, 1, len(resp_headers))
    for col in range(1, len(resp_headers) + 1):
        ws_resp.column_dimensions[get_column_letter(col)].width = 14

    out = OUT_DIR / "FPL-Manager-Survey-v1.xlsx"
    wb.save(out)
    return out


if __name__ == "__main__":
    # Prefer versioned names if originals are open/locked in the IDE
    pptx_path = build_pptx()
    xlsx_path = build_xlsx()
    print(pptx_path)
    print(xlsx_path)
