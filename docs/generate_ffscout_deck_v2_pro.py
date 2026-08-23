# -*- coding: utf-8 -*-
"""
Professional FFScout China collaboration deck (v2 evidence).
Clean corporate layout: navy / ink / forest accent, card KPIs, section structure.
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parent / "FFScout-China-Collaboration-Deck-Professional.pptx"

NAVY = RGBColor(0x0B, 0x1F, 0x33)
NAVY_MID = RGBColor(0x15, 0x32, 0x4D)
INK = RGBColor(0x1C, 0x24, 0x2E)
MUTED = RGBColor(0x5A, 0x67, 0x74)
ACCENT = RGBColor(0x0D, 0x6E, 0x56)
ACCENT_SOFT = RGBColor(0xE6, 0xF2, 0xEE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OFF = RGBColor(0xF7, 0xF8, 0xFA)
LINE = RGBColor(0xE2, 0xE8, 0xEF)
WARN = RGBColor(0x8A, 0x5A, 0x00)
FONT = "Calibri"


def run(p, text, size=14, bold=False, color=INK):
    r = p.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return r


def rect(slide, l, t, w, h, fill, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
    return s


def blank(prs, dark=False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rect(slide, 0, 0, Inches(13.333), Inches(7.5), NAVY if dark else WHITE)
    if not dark:
        rect(slide, 0, 0, Inches(13.333), Inches(0.06), ACCENT)
        rect(slide, 0, 0, Inches(0.08), Inches(7.5), NAVY)
    return slide


def footer(slide, n, total):
    rect(slide, Inches(0.7), Inches(7.05), Inches(12), Inches(0.015), LINE)
    box = slide.shapes.add_textbox(Inches(0.7), Inches(7.12), Inches(9.2), Inches(0.28))
    run(box.text_frame.paragraphs[0], f"FALEAGUE AI  ·  Confidential  ·  Scout discussion  ·  {n}/{total}", 9, False, MUTED)
    box2 = slide.shapes.add_textbox(Inches(10.2), Inches(7.12), Inches(2.5), Inches(0.28))
    box2.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
    run(box2.text_frame.paragraphs[0], "v2  ·  Survey evidence", 9, False, MUTED)


def header(slide, section, title, subtitle=None):
    box = slide.shapes.add_textbox(Inches(0.7), Inches(0.28), Inches(12), Inches(0.28))
    run(box.text_frame.paragraphs[0], section.upper(), 10, True, ACCENT)
    box2 = slide.shapes.add_textbox(Inches(0.7), Inches(0.52), Inches(12), Inches(0.5))
    run(box2.text_frame.paragraphs[0], title, 26, True, NAVY)
    rect(slide, Inches(0.7), Inches(1.08), Inches(1.15), Inches(0.04), ACCENT)
    if subtitle:
        box3 = slide.shapes.add_textbox(Inches(0.7), Inches(1.18), Inches(12), Inches(0.3))
        run(box3.text_frame.paragraphs[0], subtitle, 12, False, MUTED)
        return Inches(1.58)
    return Inches(1.32)


def bullets(slide, l, t, w, h, items, size=14):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(7)
        run(p, f"•  {item}", size, False, INK)


def kpi_cards(slide, items, top):
    n = len(items)
    gap = 0.14
    total_w = 12.0
    w = (total_w - gap * (n - 1)) / n
    height = Inches(1.12)
    for i, (val, lab) in enumerate(items):
        x = Inches(0.7 + i * (w + gap))
        rect(slide, x, top, Inches(w), height, OFF, LINE)
        rect(slide, x, top, Inches(w), Inches(0.05), ACCENT)
        vb = slide.shapes.add_textbox(x + Inches(0.1), top + Inches(0.18), Inches(w - 0.18), Inches(0.42))
        vb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        run(vb.text_frame.paragraphs[0], val, 22, True, ACCENT)
        lb = slide.shapes.add_textbox(x + Inches(0.08), top + Inches(0.62), Inches(w - 0.12), Inches(0.42))
        lb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        run(lb.text_frame.paragraphs[0], lab, 10, False, MUTED)


def callout(slide, l, t, w, h, text, tone="info"):
    fill = ACCENT_SOFT if tone == "info" else RGBColor(0xFB, 0xF3, 0xE0)
    bar = ACCENT if tone == "info" else WARN
    rect(slide, l, t, w, h, fill, LINE)
    rect(slide, l, t, Inches(0.08), h, bar)
    box = slide.shapes.add_textbox(l + Inches(0.25), t + Inches(0.14), w - Inches(0.4), h - Inches(0.22))
    box.text_frame.word_wrap = True
    run(box.text_frame.paragraphs[0], text, 12, False, NAVY)


def question_panel(slide, top, items):
    """
    Show related survey questions under the slide header.
    items: list of (qid, question_text)
    Returns y position after the panel.
    """
    n = len(items)
    # estimate height: title row + n lines
    h = Inches(0.38 + n * 0.38)
    rect(slide, Inches(0.7), top, Inches(12), h, ACCENT_SOFT, LINE)
    rect(slide, Inches(0.7), top, Inches(0.08), h, ACCENT)
    lab = slide.shapes.add_textbox(Inches(0.95), top + Inches(0.08), Inches(11.5), Inches(0.28))
    run(lab.text_frame.paragraphs[0], "Related survey questions", 10, True, ACCENT)
    for i, (qid, qtext) in enumerate(items):
        y = top + Inches(0.36 + i * 0.36)
        # qid chip
        chip_w = Inches(0.85)
        rect(slide, Inches(0.95), y, chip_w, Inches(0.28), NAVY)
        qb = slide.shapes.add_textbox(Inches(0.95), y + Inches(0.02), chip_w, Inches(0.26))
        qb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        run(qb.text_frame.paragraphs[0], qid, 10, True, WHITE)
        tb = slide.shapes.add_textbox(Inches(1.95), y + Inches(0.02), Inches(10.4), Inches(0.28))
        run(tb.text_frame.paragraphs[0], qtext, 11, False, INK)
    return top + h + Inches(0.14)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slides_meta = []  # (slide, page_num) for footer pass — footers applied immediately with known total

    TOTAL = 17
    n = 0

    # 01 Title
    n += 1
    s = blank(prs, dark=True)
    rect(s, Inches(0.7), Inches(2.05), Inches(1.5), Inches(0.05), ACCENT)
    b = s.shapes.add_textbox(Inches(0.7), Inches(2.25), Inches(11.8), Inches(0.85))
    run(b.text_frame.paragraphs[0], "FPL in China × Fantasy Football Scout", 30, True, WHITE)
    b = s.shapes.add_textbox(Inches(0.7), Inches(3.15), Inches(11.8), Inches(0.4))
    run(b.text_frame.paragraphs[0], "Collaboration discussion  ·  Evidence deck", 17, False, ACCENT)
    b = s.shapes.add_textbox(Inches(0.7), Inches(4.0), Inches(11.5), Inches(2.2))
    tf = b.text_frame
    for i, line in enumerate(
        [
            "Ray Wong  ·  FALEAGUE AI / Faleague",
            "faleague-ai.com   ·   faleague.cn",
            "Survey evidence: Wenjuanxing  ·  n=65 (WeChat)  ·  Focus n=60 (Greater China + SG/MY)",
            "Overseas brand names hidden in fielding to reduce bias",
        ]
    ):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        run(p, line, 13, False, RGBColor(0xB0, 0xBE, 0xCA))
    b = s.shapes.add_textbox(Inches(0.7), Inches(6.85), Inches(12), Inches(0.3))
    run(b.text_frame.paragraphs[0], f"Confidential discussion material  ·  {n}/{TOTAL}", 10, False, RGBColor(0x7A, 0x8A, 0x99))

    # 02 Agenda
    n += 1
    s = blank(prs)
    y0 = header(s, "Overview", "Agenda", "≈25 minutes presentation, then open discussion")
    agenda = [
        ("01", "Fit", "Why Scout × FALEAGUE is complementary"),
        ("02", "Market", "Realistic framing — no vanity China headcount"),
        ("03", "Evidence", "Survey: sources, bridge, WTP, hub demand"),
        ("04", "Offer", "Product snapshot and collaboration options"),
        ("05", "Ask", "Attributed content pilot and decision points"),
    ]
    for i, (num, short, desc) in enumerate(agenda):
        top = y0 + Inches(0.08 + i * 0.88)
        rect(s, Inches(0.7), top, Inches(12), Inches(0.78), OFF, LINE)
        rect(s, Inches(0.7), top, Inches(0.08), Inches(0.78), ACCENT)
        nb = s.shapes.add_textbox(Inches(1.0), top + Inches(0.2), Inches(0.7), Inches(0.4))
        run(nb.text_frame.paragraphs[0], num, 18, True, ACCENT)
        tb = s.shapes.add_textbox(Inches(1.8), top + Inches(0.18), Inches(2.4), Inches(0.45))
        run(tb.text_frame.paragraphs[0], short, 18, True, NAVY)
        db = s.shapes.add_textbox(Inches(4.4), top + Inches(0.22), Inches(8), Inches(0.4))
        run(db.text_frame.paragraphs[0], desc, 14, False, MUTED)
    footer(s, n, TOTAL)

    # 03 Why
    n += 1
    s = blank(prs)
    y0 = header(s, "Context", "Why we are talking")
    for i, text in enumerate(
        [
            "Fantasy Football Scout is the category authority for serious FPL managers",
            "FALEAGUE operates a live bilingual product with WeChat-native distribution",
            "We ran an unbranded survey first — this deck brings evidence, not hypotheses alone",
            "The question: does a light, attributed China channel create value for Scout?",
        ]
    ):
        top = y0 + Inches(i * 0.72)
        rect(s, Inches(0.7), top, Inches(12), Inches(0.6), OFF if i % 2 == 0 else WHITE, LINE)
        bx = s.shapes.add_textbox(Inches(0.95), top + Inches(0.15), Inches(11.5), Inches(0.35))
        run(bx.text_frame.paragraphs[0], text, 14, False, INK)
    callout(
        s,
        Inches(0.7),
        Inches(5.55),
        Inches(12),
        Inches(1.15),
        "Positioning: Scout owns trust and depth. FALEAGUE owns China access and product distribution. "
        "Survey shows Chinese managers already consume overseas signal — often without clean attribution.",
    )
    footer(s, n, TOTAL)

    # 04 Market
    n += 1
    s = blank(prs)
    y0 = header(s, "Market", "What we will and will not claim")
    rect(s, Inches(0.7), y0, Inches(5.9), Inches(4.55), OFF, LINE)
    rect(s, Inches(0.7), y0, Inches(5.9), Inches(0.06), ACCENT)
    t = s.shapes.add_textbox(Inches(0.95), y0 + Inches(0.25), Inches(5.4), Inches(0.35))
    run(t.text_frame.paragraphs[0], "Global FPL", 14, True, NAVY)
    bullets(
        s,
        Inches(0.95),
        y0 + Inches(0.75),
        Inches(5.4),
        Inches(3.5),
        [
            "~11.5M managers (2024/25)",
            "~13M+ reported (2025/26)",
            "Strong late-season engagement globally",
            "Useful as scale context only",
        ],
        13,
    )
    rect(s, Inches(6.75), y0, Inches(5.9), Inches(4.55), OFF, LINE)
    rect(s, Inches(6.75), y0, Inches(5.9), Inches(0.06), WARN)
    t = s.shapes.add_textbox(Inches(7.0), y0 + Inches(0.25), Inches(5.4), Inches(0.35))
    run(t.text_frame.paragraphs[0], "China — disciplined", 14, True, NAVY)
    bullets(
        s,
        Inches(7.0),
        y0 + Inches(0.75),
        Inches(5.4),
        Inches(3.5),
        [
            "No official China FPL census published",
            "We will not invent a headcount",
            "Proxy: China fantasy sports ~$1.2B, high-teens CAGR",
            "Survey = WeChat-cohort behaviour & WTP",
        ],
        13,
    )
    footer(s, n, TOTAL)

    # 05 Thesis
    n += 1
    s = blank(prs)
    y0 = header(s, "Thesis", "The China distribution gap")
    steps = [
        ("1", "Pull", "Managers take cues from overseas creators and specialist sites"),
        ("2", "Package", "Local creators translate / rewrite for WeChat, 小红书, short video"),
        ("3", "Gap", "Signal arrives — often delayed, partial, weakly attributed"),
        ("4", "Choice", "Stay unofficial, or build an attributed, quality-controlled channel"),
    ]
    for i, (num, title, desc) in enumerate(steps):
        top = y0 + Inches(i * 1.1)
        rect(s, Inches(0.7), top, Inches(12), Inches(0.95), OFF, LINE)
        circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.95), top + Inches(0.24), Inches(0.48), Inches(0.48))
        circ.fill.solid()
        circ.fill.fore_color.rgb = ACCENT
        circ.line.fill.background()
        nb = s.shapes.add_textbox(Inches(0.95), top + Inches(0.3), Inches(0.48), Inches(0.4))
        nb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        run(nb.text_frame.paragraphs[0], num, 14, True, WHITE)
        tb = s.shapes.add_textbox(Inches(1.7), top + Inches(0.18), Inches(2.2), Inches(0.35))
        run(tb.text_frame.paragraphs[0], title, 16, True, NAVY)
        db = s.shapes.add_textbox(Inches(4.0), top + Inches(0.22), Inches(8.4), Inches(0.5))
        run(db.text_frame.paragraphs[0], desc, 14, False, INK)
    footer(s, n, TOTAL)

    # 06 Sample
    n += 1
    s = blank(prs)
    y0 = header(s, "Evidence 01", "Who answered", "Unbranded Wenjuanxing survey · Aug 2026")
    kpi_cards(s, [("65", "Responses"), ("60", "CN / HK / TW / SG-MY"), ("100%", "WeChat recruit"), ("78%", "Bilingual")], y0)
    rows = [
        ("Playing status", "This season 66.7%  ·  Previously 23.3%  ·  Not started yet 10%"),
        ("Experience", "4–6 seasons 41.7%  ·  2–3 35%  ·  7+ 18.3%"),
        ("Rank band", "Top 100k 41.7%  ·  Top 1M 38.3%  ·  Top 10k 15%"),
        ("Demographics", "Age 25–34 58%  ·  35–44 32%"),
    ]
    for i, (lab, val) in enumerate(rows):
        top = y0 + Inches(1.35 + i * 0.7)
        rect(s, Inches(0.7), top, Inches(12), Inches(0.58), WHITE, LINE)
        a = s.shapes.add_textbox(Inches(0.95), top + Inches(0.14), Inches(2.6), Inches(0.35))
        run(a.text_frame.paragraphs[0], lab, 13, True, NAVY)
        b = s.shapes.add_textbox(Inches(3.7), top + Inches(0.14), Inches(8.7), Inches(0.35))
        run(b.text_frame.paragraphs[0], val, 13, False, INK)
    callout(
        s,
        Inches(0.7),
        Inches(6.3),
        Inches(12),
        Inches(0.55),
        "Read as directional evidence from a WeChat convenience sample — not a national FPL census.",
        "warn",
    )
    footer(s, n, TOTAL)

    # 07 Sources
    n += 1
    s = blank(prs)
    y0 = header(s, "Evidence 02", "Information sources", "Focus n=60")
    kpi_cards(
        s,
        [("66.7%", "Overseas site/creator"), ("70%", "CN creators / 公众号"), ("63%", "WeChat as source"), ("91.7%", "WeChat group habit")],
        y0,
    )
    headers = ["Channel", "Use regularly", "In top-3 importance"]
    data = [
        ("Overseas creators", "58%", "53%"),
        ("CN creators / 公众号", "70%", "50%"),
        ("Official FPL", "58%", "38%"),
        ("WeChat", "63%", "37%"),
        ("Overseas sites / tools", "37%", "28%"),
        ("Local translated posts", "37%", "25%"),
    ]
    top = y0 + Inches(1.35)
    rect(s, Inches(0.7), top, Inches(12), Inches(0.4), NAVY)
    widths = [4.5, 3.7, 3.8]
    x = 0.7
    for i, h in enumerate(headers):
        bx = s.shapes.add_textbox(Inches(x + 0.15), top + Inches(0.08), Inches(widths[i] - 0.2), Inches(0.28))
        run(bx.text_frame.paragraphs[0], h, 12, True, WHITE)
        x += widths[i]
    for ri, row in enumerate(data):
        rt = top + Inches(0.4 + ri * 0.52)
        rect(s, Inches(0.7), rt, Inches(12), Inches(0.52), OFF if ri % 2 == 0 else WHITE, LINE)
        x = 0.7
        for ci, cell in enumerate(row):
            bx = s.shapes.add_textbox(Inches(x + 0.15), rt + Inches(0.12), Inches(widths[ci] - 0.2), Inches(0.3))
            run(bx.text_frame.paragraphs[0], cell, 12, ci > 0, ACCENT if ci > 0 else INK)
            x += widths[ci]
    footer(s, n, TOTAL)

    # 08 Bridge
    n += 1
    s = blank(prs)
    y0 = header(s, "Evidence 03", "The translation bridge — validated")
    kpi_cards(
        s,
        [("46.7%", "Translated tips weekly+"), ("78.3%", "Use posts as input"), ("31.7%", "Rely on CN repost"), ("63.3%", "Also read EN original")],
        y0,
    )
    findings = [
        ("Consumption path", "EN original 63%  ·  Self-translate 38%  ·  CN repost 32%  ·  Rarely use overseas 10%"),
        ("EN analysis need", "Need translation / summary 40%  ·  Rely heavily 25%  ·  Mostly EN 22%"),
        ("When they see posts", "One input among several 75% — not blind copying"),
        ("Barrier to overseas tools", "None 33%  ·  Access/VPN 20%  ·  Paywall 17%  ·  Unknown 15%  ·  Language 12%"),
    ]
    for i, (lab, val) in enumerate(findings):
        top = y0 + Inches(1.4 + i * 0.85)
        rect(s, Inches(0.7), top, Inches(12), Inches(0.72), OFF, LINE)
        rect(s, Inches(0.7), top, Inches(0.08), Inches(0.72), ACCENT)
        a = s.shapes.add_textbox(Inches(1.0), top + Inches(0.2), Inches(3.1), Inches(0.35))
        run(a.text_frame.paragraphs[0], lab, 13, True, NAVY)
        b = s.shapes.add_textbox(Inches(4.2), top + Inches(0.2), Inches(8.2), Inches(0.4))
        run(b.text_frame.paragraphs[0], val, 12, False, INK)
    footer(s, n, TOTAL)

    # 09 Jobs + WTP
    n += 1
    s = blank(prs)
    y0 = header(s, "Evidence 04", "Jobs, formats & willingness to pay")
    cols = [
        ("Hardest decisions", ["Transfers — 67%", "Bench order — 40%", "Captain — 38%", "Differentials — 25%"]),
        ("Preferred CN formats", ["Short daily card — 58%", "Long analysis — 45%", "Data tables — 32%", "Short video — 17%"]),
        ("WTP (CNY / month)", ["Won’t pay — 55%", "< ¥20 — 30%", "¥20–40 — 8%", "Prefer annual — 7%"]),
    ]
    for i, (title, items) in enumerate(cols):
        x = Inches(0.7 + i * 4.15)
        rect(s, x, y0, Inches(3.95), Inches(3.45), OFF, LINE)
        rect(s, x, y0, Inches(3.95), Inches(0.06), ACCENT)
        tb = s.shapes.add_textbox(x + Inches(0.2), y0 + Inches(0.25), Inches(3.5), Inches(0.35))
        run(tb.text_frame.paragraphs[0], title, 13, True, NAVY)
        bullets(s, x + Inches(0.2), y0 + Inches(0.75), Inches(3.5), Inches(2.5), items, 13)
    callout(
        s,
        Inches(0.7),
        Inches(5.7),
        Inches(12),
        Inches(1.0),
        "Commercial read: any paid intent 45%; ¥20+/mo or annual only 15%. "
        "Lead with free attributed content; keep China pricing careful. Top paid-tool intents: squad-tied AI 28%, transfer planner 20%.",
    )
    footer(s, n, TOTAL)

    # 10 Hub
    n += 1
    s = blank(prs)
    y0 = header(s, "Evidence 05", "Demand for an attributed Chinese hub")
    kpi_cards(
        s,
        [("80%", "Interested in CN hub"), ("35%", "Follow or consider paid"), ("27%", "Quality-first"), ("60%", "Need depth ≥ original")],
        y0,
    )
    rect(s, Inches(0.7), y0 + Inches(1.35), Inches(12), Inches(4.0), OFF, LINE)
    t = s.shapes.add_textbox(Inches(0.95), y0 + Inches(1.55), Inches(11.5), Inches(0.35))
    run(t.text_frame.paragraphs[0], "Trust conditions for a partnership (Q25)", 14, True, NAVY)
    trust = [
        ("Depth ≥ overseas original", 60),
        ("Stable weekly cadence", 47),
        ("High-quality translation", 47),
        ("Clear source attribution", 43),
        ("Fair China pricing", 25),
    ]
    for i, (lab, pct) in enumerate(trust):
        top = y0 + Inches(2.1 + i * 0.55)
        a = s.shapes.add_textbox(Inches(0.95), top, Inches(4.2), Inches(0.35))
        run(a.text_frame.paragraphs[0], lab, 12, False, INK)
        track_w = 6.2
        rect(s, Inches(5.3), top + Inches(0.08), Inches(track_w), Inches(0.22), WHITE, LINE)
        rect(s, Inches(5.3), top + Inches(0.08), Inches(track_w * pct / 100), Inches(0.22), ACCENT)
        pctb = s.shapes.add_textbox(Inches(11.6), top, Inches(0.9), Inches(0.35))
        run(pctb.text_frame.paragraphs[0], f"{pct}%", 12, True, ACCENT)
    footer(s, n, TOTAL)

    # 11 Unaided
    n += 1
    s = blank(prs)
    y0 = header(s, "Evidence 06", "Unaided trust names", "Q7 open answers · n=34 usable · brands were not prompted")
    bullets(
        s,
        Inches(0.7),
        y0,
        Inches(12),
        Inches(3.6),
        [
            "Creators dominate recall: FPL Harry ~32% of answered opens",
            "Also mentioned: All About Fantasy, FPL Focal, FPL Gameweek, Andy, Big Man, Sam",
            "Fantasy Football Scout / FFS appears in 3 answers (~9% of answered)",
            "Interpretation: Scout has unaided presence for some managers — not a mass awareness claim",
            "Opportunity: be the attributed research brand behind Chinese packaging",
        ],
        14,
    )
    callout(
        s,
        Inches(0.7),
        Inches(5.55),
        Inches(12),
        Inches(1.1),
        "Recommended language: “Scout already appears in some unaided trust lists.” "
        "Do not state a national awareness percentage from this sample.",
        "warn",
    )
    footer(s, n, TOTAL)

    # 12 Product
    n += 1
    s = blank(prs)
    y0 = header(s, "Product", "What FALEAGUE has built")
    bullets(
        s,
        Inches(0.7),
        y0,
        Inches(12),
        Inches(4.8),
        [
            "Live hub: faleague-ai.com (EN + ZH)  ·  sister community faleague.cn",
            "Tool-grounded AI analyst + Insights research tables + Entry ID planner",
            "WeChat pipelines: daily cards, Best of Position hooks, 小程序 shell",
            "Survey alignment: respondents want short daily cards + tables — matches our distribution",
            "We already cite Scout in Chinese digests — amplifying quality signal into WeChat",
            "Not claiming Opta parity — China-facing product and distribution layer",
        ],
        14,
    )
    footer(s, n, TOTAL)

    # 13 Options
    n += 1
    s = blank(prs)
    y0 = header(s, "Collaboration", "Options — updated by evidence")
    opts = [
        ("A  ·  Recommended next", "Content & attribution", "Scout-approved Chinese WeChat / XHS cards with clear source and link-back"),
        ("B", "Co-branded Insights", "Scout editorial judgement + FALEAGUE bilingual Insights UI"),
        ("C", "Membership funnel", "Soft CTA to Scout membership for advanced bilingual managers; unique links"),
        ("D  ·  Complete", "Joint research", "Unbranded survey done; optional later co-published China-lens note"),
    ]
    for i, (tag, title, desc) in enumerate(opts):
        top = y0 + Inches(i * 1.1)
        rect(s, Inches(0.7), top, Inches(12), Inches(0.95), OFF, LINE)
        rect(s, Inches(0.7), top, Inches(0.08), Inches(0.95), ACCENT if i == 0 else LINE)
        tg = s.shapes.add_textbox(Inches(1.0), top + Inches(0.15), Inches(3.4), Inches(0.3))
        run(tg.text_frame.paragraphs[0], tag, 11, True, ACCENT)
        ti = s.shapes.add_textbox(Inches(4.5), top + Inches(0.12), Inches(7.9), Inches(0.3))
        run(ti.text_frame.paragraphs[0], title, 15, True, NAVY)
        de = s.shapes.add_textbox(Inches(4.5), top + Inches(0.5), Inches(7.9), Inches(0.35))
        run(de.text_frame.paragraphs[0], desc, 12, False, MUTED)
    footer(s, n, TOTAL)

    # 14 Pilot
    n += 1
    s = blank(prs)
    y0 = header(s, "Plan", "Proposed next 90 days")
    pilot = [
        ("Done", "Unbranded survey fielded and analysed", "This evidence deck"),
        ("Weeks 1–2", "Align attribution rules; pick 2–3 safe content types", "Written one-pager"),
        ("Weeks 2–6", "Ship attributed WeChat / XHS cards with unique links", "CTR / deep-link opens"),
        ("Weeks 6–10", "Optional co-branded Insights asset for a key GW", "Engagement + qualitative"),
        ("Weeks 10–12", "Review: expand, membership test, or stop", "Go / no-go"),
    ]
    rect(s, Inches(0.7), y0, Inches(12), Inches(0.4), NAVY)
    for x, h, w in [(0.9, "When", 2.2), (3.2, "Action", 6.2), (9.5, "Success signal", 3.0)]:
        bx = s.shapes.add_textbox(Inches(x), y0 + Inches(0.08), Inches(w), Inches(0.28))
        run(bx.text_frame.paragraphs[0], h, 11, True, WHITE)
    for i, (when, action, signal) in enumerate(pilot):
        top = y0 + Inches(0.4 + i * 0.78)
        rect(s, Inches(0.7), top, Inches(12), Inches(0.78), OFF if i % 2 == 0 else WHITE, LINE)
        for x, text, w, bold in [(0.9, when, 2.2, True), (3.2, action, 6.2, False), (9.5, signal, 3.0, False)]:
            bx = s.shapes.add_textbox(Inches(x), top + Inches(0.22), Inches(w), Inches(0.4))
            run(bx.text_frame.paragraphs[0], text, 12, bold, ACCENT if bold else INK)
    footer(s, n, TOTAL)

    # 15 Ask
    n += 1
    s = blank(prs)
    y0 = header(s, "Ask", "What we need from this conversation")
    asks = [
        "Interest in a time-boxed attributed content pilot (Option A)",
        "A single point of contact for brand and content approval",
        "Clarity on what must never be localised without editorial control",
        "Whether brand reach alone is the goal, or membership / licensing later",
        "Optional: comfort citing this survey as joint directional research",
    ]
    for i, text in enumerate(asks):
        top = y0 + Inches(i * 0.85)
        rect(s, Inches(0.7), top, Inches(12), Inches(0.72), OFF, LINE)
        num = s.shapes.add_textbox(Inches(0.95), top + Inches(0.2), Inches(0.55), Inches(0.35))
        run(num.text_frame.paragraphs[0], f"{i+1:02d}", 14, True, ACCENT)
        bx = s.shapes.add_textbox(Inches(1.65), top + Inches(0.2), Inches(10.7), Inches(0.4))
        run(bx.text_frame.paragraphs[0], text, 14, False, INK)
    footer(s, n, TOTAL)

    # 16 Discussion
    n += 1
    s = blank(prs)
    y0 = header(s, "Discussion", "Prompts")
    prompts = [
        "Is China a strategic priority for Scout over the next 12–24 months, or opportunistic?",
        "Does 80% hub interest plus WeChat distribution change that priority?",
        "How should Scout show up relative to the creator layer (Harry et al. dominate unaided recall)?",
        "What commercial path is acceptable given ¥ price sensitivity?",
        "Who would own the pilot internally if we proceed?",
    ]
    for i, text in enumerate(prompts):
        top = y0 + Inches(i * 0.9)
        rect(s, Inches(0.7), top, Inches(12), Inches(0.78), WHITE, LINE)
        rect(s, Inches(0.7), top, Inches(0.08), Inches(0.78), ACCENT)
        bx = s.shapes.add_textbox(Inches(1.05), top + Inches(0.2), Inches(11.3), Inches(0.45))
        run(bx.text_frame.paragraphs[0], text, 14, False, INK)
    footer(s, n, TOTAL)

    # 17 Close
    n += 1
    s = blank(prs, dark=True)
    rect(s, Inches(0.7), Inches(1.55), Inches(1.5), Inches(0.05), ACCENT)
    b = s.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(11.8), Inches(0.5))
    run(b.text_frame.paragraphs[0], "Closing", 28, True, WHITE)
    b = s.shapes.add_textbox(Inches(0.7), Inches(2.6), Inches(11.8), Inches(2.2))
    tf = b.text_frame
    for i, line in enumerate(
        [
            "Scout = trust and depth  ·  FALEAGUE = bilingual product + WeChat China distribution",
            "Evidence: overseas signal is consumed; WeChat is the habit loop; hub demand is real",
            "Ask: a reversible attributed content pilot — not a big-bang JV",
        ]
    ):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(14)
        run(p, f"•  {line}", 15, False, RGBColor(0xD0, 0xDA, 0xE2))
    rect(s, Inches(0.7), Inches(5.15), Inches(12), Inches(1.45), NAVY_MID)
    eb = s.shapes.add_textbox(Inches(0.95), Inches(5.35), Inches(11.5), Inches(1.15))
    tf = eb.text_frame
    run(tf.paragraphs[0], "Evidence one-liner", 11, True, ACCENT)
    p2 = tf.add_paragraph()
    p2.space_before = Pt(8)
    run(
        p2,
        "In this WeChat sample (focus n=60), 67% use overseas FPL sources, 92% live in WeChat groups, "
        "and 80% want an attributed Chinese hub — pricing must stay careful in ¥.",
        13,
        False,
        WHITE,
    )
    b = s.shapes.add_textbox(Inches(0.7), Inches(6.85), Inches(12), Inches(0.3))
    run(b.text_frame.paragraphs[0], f"Thank you  ·  {n}/{TOTAL}", 11, False, RGBColor(0x8A, 0x9A, 0xA8))

    prs.save(OUT)
    return OUT, n


if __name__ == "__main__":
    path, count = build()
    print(path)
    print("slides", count)
