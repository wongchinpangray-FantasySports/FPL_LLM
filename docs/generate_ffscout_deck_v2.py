# -*- coding: utf-8 -*-
"""Generate FFScout China collaboration deck v2 with survey evidence merged."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

OUT_DIR = Path(__file__).resolve().parent

NAVY = RGBColor(0x0F, 0x1C, 0x2E)
INK = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x5C, 0x67, 0x73)
ACCENT = RGBColor(0x0B, 0x6E, 0x4F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF4, 0xF6, 0xF8)
LINE = RGBColor(0xD0, 0xD7, 0xDE)


def set_run(run, text, size=18, bold=False, color=INK, font_name="Calibri"):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name


def clear_and_write(tf, lines, space_after=6):
    tf.clear()
    tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, tuple):
            text, is_bold, sz = item[0], item[1] if len(item) > 1 else False, item[2] if len(item) > 2 else 14
            col = item[3] if len(item) > 3 else INK
        else:
            text, is_bold, sz, col = item, False, 14, INK
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(space_after)
        run = p.add_run()
        set_run(run, text, sz, is_bold, col)


def add_bullets(slide, left, top, width, height, bullets, size=15):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        run = p.add_run()
        set_run(run, f"•  {b}", size, False, INK)
    return box


def add_footer(slide, page, total):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(12), Inches(0.3))
    p = box.text_frame.paragraphs[0]
    set_run(
        p.add_run(),
        f"FALEAGUE AI  ·  Scout conversation deck v2 (survey evidence)  ·  {page}/{total}",
        10,
        False,
        MUTED,
    )


def add_accent_bar(slide):
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.12), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()


def blank_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()
    add_accent_bar(slide)
    return slide


def title_block(slide, title, subtitle=None):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.32), Inches(12.2), Inches(0.65))
    set_run(box.text_frame.paragraphs[0].add_run(), title, 26, True, NAVY)
    if subtitle:
        box2 = slide.shapes.add_textbox(Inches(0.55), Inches(0.92), Inches(12.2), Inches(0.35))
        set_run(box2.text_frame.paragraphs[0].add_run(), subtitle, 13, False, MUTED)


def kpi_row(slide, items, top=Inches(1.4)):
    """items: list of (value, label)"""
    n = len(items)
    width = 12.2 / n
    for i, (val, lab) in enumerate(items):
        x = Inches(0.55 + i * width)
        shape = slide.shapes.add_shape(1, x, top, Inches(width - 0.12), Inches(1.05))
        shape.fill.solid()
        shape.fill.fore_color.rgb = LIGHT
        shape.line.color.rgb = LINE
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        set_run(p.add_run(), val, 22, True, ACCENT)
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        set_run(p2.add_run(), lab, 11, False, MUTED)


def build_pptx():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    total = 18

    # 1 Title
    slide = blank_slide(prs)
    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.9), Inches(11.5), Inches(1.1))
    set_run(box.text_frame.paragraphs[0].add_run(), "FPL in China × Fantasy Football Scout", 34, True, NAVY)
    box2 = slide.shapes.add_textbox(Inches(0.8), Inches(3.1), Inches(11.5), Inches(0.45))
    set_run(box2.text_frame.paragraphs[0].add_run(), "Collaboration deck v2 — with survey evidence", 20, False, ACCENT)
    add_bullets(
        slide,
        Inches(0.8),
        Inches(3.9),
        Inches(11),
        Inches(2.2),
        [
            "Presenter: Ray Wong · FALEAGUE AI / Faleague",
            "Product: faleague-ai.com  ·  Sister: faleague.cn",
            "Survey: Wenjuanxing · n=65 (WeChat) · Focus Greater China + SG/MY n=60",
            "Brand names were hidden in fielding to reduce bias",
        ],
        15,
    )
    add_footer(slide, 1, total)

    # 2 Agenda
    slide = blank_slide(prs)
    title_block(slide, "Agenda", "≈25 min + discussion")
    add_bullets(
        slide,
        Inches(0.7),
        Inches(1.5),
        Inches(12),
        Inches(5),
        [
            "Mutual fit — Scout trust/depth × FALEAGUE China distribution",
            "Market reality (no invented China FPL headcount)",
            "Survey evidence: sources, translation bridge, WTP, hub interest",
            "Product snapshot + collaboration options",
            "Proposed next step: attributed content pilot",
        ],
        17,
    )
    add_footer(slide, 2, total)

    # 3 Why
    slide = blank_slide(prs)
    title_block(slide, "1 · Why we’re talking")
    add_bullets(
        slide,
        Inches(0.7),
        Inches(1.5),
        Inches(12),
        Inches(3.5),
        [
            "Scout is the category authority for serious FPL managers",
            "We built a live bilingual product with WeChat-native distribution",
            "We fielded an unbranded survey first — this deck returns with evidence",
            "Ask: does a light, attributed China channel create value for Scout?",
        ],
        16,
    )
    box = slide.shapes.add_shape(1, Inches(0.7), Inches(5.2), Inches(11.8), Inches(1.1))
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT
    box.line.color.rgb = LINE
    clear_and_write(
        box.text_frame,
        [
            (
                "One line: Scout owns trust and depth. We own China access and product distribution. Survey shows the bridge already exists — mostly unofficial.",
                True,
                15,
                NAVY,
            )
        ],
    )
    add_footer(slide, 3, total)

    # 4 Market
    slide = blank_slide(prs)
    title_block(slide, "2 · Market reality (realistic framing)")
    add_bullets(
        slide,
        Inches(0.7),
        Inches(1.5),
        Inches(12),
        Inches(5),
        [
            "Global FPL: ~11.5M (24/25) → ~13M+ reported (25/26) — order of magnitude",
            "Official China FPL manager census is NOT published — we will not invent one",
            "Proxy only: China fantasy sports ~USD 1.2B (2024), high-teens CAGR (industry reports) — not FPL revenue",
            "PL fandom in China is mass historically; FPL conversion is the open question",
            "This survey sizes behaviour & WTP in a WeChat-recruited cohort — not national population",
        ],
        15,
    )
    add_footer(slide, 4, total)

    # 5 China gap thesis
    slide = blank_slide(prs)
    title_block(slide, "3 · The China gap (thesis we tested)")
    add_bullets(
        slide,
        Inches(0.7),
        Inches(1.5),
        Inches(12),
        Inches(5),
        [
            "Chinese FPLers pull recommendations from overseas creators/sites",
            "Local creators often translate/rewrite that content onto WeChat / 小红书 / short video",
            "Signal already reaches China — often with weak or missing brand attribution",
            "Strategic question for Scout: stay unofficial & uncontrolled, or become attributed + quality-controlled?",
            "Survey was designed to validate this bridge without naming partner brands",
        ],
        15,
    )
    add_footer(slide, 5, total)

    # 6 Survey caveat + sample
    slide = blank_slide(prs)
    title_block(
        slide,
        "4 · Survey evidence — sample",
        "Wenjuanxing · Aug 2026 · unbranded fielding",
    )
    kpi_row(
        slide,
        [
            ("65", "Total responses"),
            ("60", "Greater China + SG/MY"),
            ("100%", "Recruited via WeChat"),
            ("78%", "Bilingual content prefer"),
        ],
    )
    add_bullets(
        slide,
        Inches(0.7),
        Inches(2.7),
        Inches(12),
        Inches(3.8),
        [
            "Playing this season 66.7% · used to play 23.3%",
            "Tenure: 4–6 seasons 41.7% · 2–3 35% · 7+ 18.3% (experienced sample)",
            "Rank band: Top 100k 41.7% · Top 1M 38.3% · Top 10k 15%",
            "Age: 25–34 58.3% · 35–44 31.7%",
            "CAVEAT: convenience WeChat sample — directional behaviour evidence, NOT a China FPL census",
        ],
        14,
    )
    add_footer(slide, 6, total)

    # 7 Sources
    slide = blank_slide(prs)
    title_block(slide, "5 · How they get FPL information", "Focus n=60 · multi-select")
    kpi_row(
        slide,
        [
            ("66.7%", "Overseas site or creator"),
            ("70%", "CN creators / 公众号"),
            ("63.3%", "WeChat as source"),
            ("91.7%", "Discuss in WeChat groups"),
        ],
    )
    add_bullets(
        slide,
        Inches(0.7),
        Inches(2.7),
        Inches(12),
        Inches(3.8),
        [
            "Regular sources: CN KOLs 70% · WeChat 63% · Official FPL 58% · Overseas creators 58%",
            "Overseas sites/tools 36.7% · Local translated/rewritten posts 36.7% · 小红书 26.7%",
            "Most important (top-3): Overseas creators 53.3% · CN creators 50% · Official 38% · WeChat 37%",
            "Implication: Scout-quality signal is valued; WeChat + CN packaging is the daily layer",
        ],
        14,
    )
    add_footer(slide, 7, total)

    # 8 Translation bridge
    slide = blank_slide(prs)
    title_block(slide, "6 · The translation bridge — validated")
    kpi_row(
        slide,
        [
            ("46.7%", "See translated tips weekly+"),
            ("78.3%", "Use those posts as input"),
            ("31.7%", "Rely on CN repost path"),
            ("63.3%", "Also read EN original"),
        ],
    )
    add_bullets(
        slide,
        Inches(0.7),
        Inches(2.7),
        Inches(12),
        Inches(3.8),
        [
            "Q8 path: EN original 63% · self-translate 38% · CN creator/repost 32% · rarely use overseas 10%",
            "Q9: need translation/summary 40% · rely heavily on EN 25% · mostly EN (skip) 22%",
            "Q16 when they see translated posts: one input among several 75% (not blind copy)",
            "Barriers to overseas tools: no barrier 33% · VPN/access 20% · paywall 17% · don’t know good ones 15% · language 12%",
            "Story: the unofficial bridge exists; attributed packaging is the upgrade",
        ],
        13,
    )
    add_footer(slide, 8, total)

    # 9 Jobs + formats + WTP
    slide = blank_slide(prs)
    title_block(slide, "7 · Jobs, formats & willingness to pay (¥)")
    add_bullets(
        slide,
        Inches(0.55),
        Inches(1.4),
        Inches(6.0),
        Inches(5),
        [
            "Hardest GW decisions: Transfers 67% · Bench 40% · Captain 38%",
            "Preferred CN formats: Short daily card 58% · Long analysis 45% · Interactive tables 32%",
            "Tools they’d pay for: Free-only 62% · Squad-tied AI 28% · Transfer planner 20%",
            "WeChat is the habit loop — short cards fit distribution",
        ],
        13,
    )
    add_bullets(
        slide,
        Inches(6.8),
        Inches(1.4),
        Inches(5.8),
        Inches(5),
        [
            "WTP (Q12): Won’t pay 55% · <¥20/mo 30% · ¥20–40 8% · Prefer annual 7%",
            "Any paid intent: 45%",
            "¥20+/mo or annual: only 15%",
            "Commercial read: start with free attributed content + light Pro; China pricing must be careful",
        ],
        13,
    )
    add_footer(slide, 9, total)

    # 10 Hub / partnership
    slide = blank_slide(prs)
    title_block(slide, "8 · Interest in attributed China hub / partnership")
    kpi_row(
        slide,
        [
            ("80%", "Want attributed CN hub"),
            ("35%", "Follow or consider paid"),
            ("26.7%", "Quality-first"),
            ("60%", "Need depth ≥ original"),
        ],
    )
    add_bullets(
        slide,
        Inches(0.7),
        Inches(2.7),
        Inches(12),
        Inches(3.8),
        [
            "Q19 hub interest: somewhat 47% · very interested 33% · neutral 18% · not interested 2%",
            "Q24 if overseas research brand partnered for CN content/tools: free-only 37% · quality-first 27% · follow 20% · consider paid 15%",
            "Trust conditions (Q25): depth ≥ overseas 60% · weekly cadence 47% · quality translation 47% · clear attribution 43%",
            "Also: fair China pricing 25% · local community 20% · no gambling framing 18%",
        ],
        13,
    )
    add_footer(slide, 10, total)

    # 11 Unaided brands
    slide = blank_slide(prs)
    title_block(
        slide,
        "9 · Unaided trust names (Q7)",
        "n=34 usable open answers · brands were NOT prompted",
    )
    add_bullets(
        slide,
        Inches(0.7),
        Inches(1.55),
        Inches(12),
        Inches(4.8),
        [
            "Creators dominate recall: FPL Harry ~32% of answered · All About Fantasy ~12%",
            "Also: FPL Focal / FPL Gameweek / Andy / Big Man / Sam (~9% each of answered)",
            "Fantasy Football Scout / FFS appears in 3 answers (~9% of answered)",
            "Use carefully: Scout has unaided presence for some managers — not a high awareness claim",
            "Creators are the top-of-mind layer; Scout can win as the attributed research brand behind CN packaging",
        ],
        15,
    )
    add_footer(slide, 11, total)

    # 12 Product
    slide = blank_slide(prs)
    title_block(slide, "10 · What we built — FALEAGUE AI")
    add_bullets(
        slide,
        Inches(0.7),
        Inches(1.5),
        Inches(12),
        Inches(5),
        [
            "Live: faleague-ai.com (EN + ZH) · sister faleague.cn",
            "Tool-grounded AI + Insights tables + Entry ID planner · WeChat daily cards / 小程序",
            "Survey product signal (secondary): 75% tried at least once; NPS ~58 among users (n=45)",
            "Most useful cited: Insights + squad builder; formats they want match our daily card + tables",
            "We already cite Scout in Chinese digests — amplifying quality signal into WeChat",
        ],
        15,
    )
    add_footer(slide, 12, total)

    # 13 Demo
    slide = blank_slide(prs)
    title_block(slide, "11 · Demo (3 minutes live)")
    add_bullets(
        slide,
        Inches(0.7),
        Inches(1.5),
        Inches(12),
        Inches(5),
        [
            "/zh/fpl/insights — Best of Position, pre-season, set-pieces, fixtures",
            "/zh/chat — tool-grounded AI with Entry ID",
            "One WeChat daily card / BOP hook with deep links",
            "Not claiming Opta parity — China product + distribution layer",
        ],
        16,
    )
    add_footer(slide, 13, total)

    # 14 Collaboration menu
    slide = blank_slide(prs)
    title_block(slide, "12 · Collaboration options (updated by evidence)")
    add_bullets(
        slide,
        Inches(0.7),
        Inches(1.5),
        Inches(12),
        Inches(5),
        [
            "A. Content & attribution (recommended next) — Scout-approved Chinese WeChat/XHS cards; clear source + link",
            "B. Co-branded China Insights page — Scout editorial + our ZH UI",
            "C. Soft membership funnel — unique links for advanced bilingual managers",
            "D. Survey — DONE (unbranded). Optional later: co-published China-lens report",
            "Evidence says start with A: short daily cards + attribution + depth; keep ¥ pricing light",
        ],
        14,
    )
    add_footer(slide, 14, total)

    # 15 Pilot
    slide = blank_slide(prs)
    title_block(slide, "13 · Proposed next 90 days")
    pilot = [
        ("Done", "Unbranded survey fielded & analysed (n=65 / focus 60)", "Evidence in this deck"),
        ("Week 1–2", "Align Scout brand/attribution rules; pick 2–3 safe content types", "One-pager"),
        ("Week 2–6", "Ship attributed WeChat/XHS cards + unique links", "CTR / deep-links"),
        ("Week 6–10", "Optional: one co-branded Insights asset for a key GW", "Engagement + qualitative"),
        ("Week 10–12", "Review: expand, membership test, or stop", "Go / no-go"),
    ]
    for i, (w, a, s) in enumerate(pilot):
        y = Inches(1.4 + i * 0.72)
        shape = slide.shapes.add_shape(1, Inches(0.55), y, Inches(12.2), Inches(0.64))
        shape.fill.solid()
        shape.fill.fore_color.rgb = LIGHT if i % 2 == 0 else WHITE
        shape.line.color.rgb = LINE
        clear_and_write(shape.text_frame, [(f"{w}  ·  {a}  →  {s}", False, 13, INK)])
    add_footer(slide, 15, total)

    # 16 Ask
    slide = blank_slide(prs)
    title_block(slide, "14 · Ask on this call")
    add_bullets(
        slide,
        Inches(0.7),
        Inches(1.55),
        Inches(12),
        Inches(5),
        [
            "Interest in a time-boxed attributed content pilot (option A)",
            "One point of contact for brand/content approval",
            "What must never be localised without editorial control?",
            "Prefer brand reach only, or also explore membership / licensing later?",
            "Optional: comfort with citing this survey as joint directional research",
        ],
        16,
    )
    add_footer(slide, 16, total)

    # 17 Discussion
    slide = blank_slide(prs)
    title_block(slide, "15 · Discussion prompts")
    add_bullets(
        slide,
        Inches(0.7),
        Inches(1.55),
        Inches(12),
        Inches(5),
        [
            "Is China a strategic priority for Scout in 12–24 months, or opportunistic?",
            "Does 80% hub interest + WeChat distribution change the priority?",
            "How should Scout show up vs creator ecosystem (Harry etc. dominate unaided recall)?",
            "What commercial path is acceptable in China given ¥ price sensitivity?",
            "Who owns the pilot internally if we proceed?",
        ],
        16,
    )
    add_footer(slide, 17, total)

    # 18 Close
    slide = blank_slide(prs)
    title_block(slide, "16 · Closing")
    add_bullets(
        slide,
        Inches(0.7),
        Inches(1.5),
        Inches(12),
        Inches(2.8),
        [
            "Scout = trust + depth · FALEAGUE = bilingual product + WeChat China distribution",
            "Survey: overseas signal is already consumed; WeChat is the habit loop; hub demand is real",
            "Ask = reversible attributed content pilot — not a big-bang JV",
        ],
        16,
    )
    box = slide.shapes.add_shape(1, Inches(0.7), Inches(4.6), Inches(11.8), Inches(1.6))
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT
    box.line.color.rgb = LINE
    clear_and_write(
        box.text_frame,
        [
            ("Evidence one-liner", True, 14, ACCENT),
            (
                "In this WeChat sample (n=60 focus), 67% use overseas FPL sources, 92% live in WeChat groups, 80% want an attributed Chinese hub — and pricing must stay careful in ¥.",
                False,
                14,
                NAVY,
            ),
        ],
    )
    add_footer(slide, 18, total)

    out = OUT_DIR / "FFScout-China-Collaboration-Deck-v2.pptx"
    prs.save(out)
    return out


if __name__ == "__main__":
    path = build_pptx()
    print(path)
